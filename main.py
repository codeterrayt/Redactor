import os
import re
import json
import spacy
import openpyxl
from pptx import Presentation
from rapidfuzz import process, fuzz
from cleanco import basename

# --- CONFIGURATION ---
# Add names here to target specific clients (e.g., ["google", "amazon"])
# Leave empty [] to auto-detect all companies found by the AI.
GLOBAL_CLIENT_LIST = [] 

# Load the Transformer model
print("Loading Deep Intelligence Model (spaCy TRF)...")
nlp = spacy.load("en_core_web_trf")

class GlobalRedactor:
    def __init__(self, target_clients=None, mapping_file="client_mapping.json"):
        self.mapping_file = mapping_file
        self.target_clients = [self.normalize(c) for c in (target_clients or []) if c]
        
        print("\n" + "="*70)
        if not self.target_clients:
            print("MODE: Auto-Discovery (Redacting all detected organizations)")
        else:
            print(f"MODE: Targeted Redaction")
            print(f"TARGETS: {', '.join(target_clients)}")
        print("="*70 + "\n")

        if os.path.exists(mapping_file):
            with open(mapping_file, 'r') as f:
                data = json.load(f)
                self.client_map = data.get('map', {})
                self.counter = data.get('counter', 1)
        else:
            self.client_map = {} 
            self.counter = 1

    def normalize(self, text):
        if not text: return ""
        base = basename(str(text))
        base = base.lower().strip()
        base = re.sub(r'[^\w\s]', '', base) 
        return base.strip()

    def is_valid_target(self, detected_name):
        if not self.target_clients:
            return True
        norm_detected = self.normalize(detected_name)
        match = process.extractOne(
            norm_detected, 
            self.target_clients, 
            scorer=fuzz.WRatio, 
            score_cutoff=90
        )
        return True if match else False

    def get_client_id(self, raw_name):
        if not self.is_valid_target(raw_name):
            return None
        norm_name = self.normalize(raw_name)
        if not norm_name or len(norm_name) < 2: return None

        if norm_name in self.client_map:
            return self.client_map[norm_name]

        sorted_keys = sorted(self.client_map.keys(), key=len, reverse=True)
        for existing_norm in sorted_keys:
            if existing_norm in norm_name or norm_name in existing_norm:
                target_id = self.client_map[existing_norm]
                self.client_map[norm_name] = target_id
                return target_id

        client_id = f"[client{self.counter}]"
        self.client_map[norm_name] = client_id
        self.counter += 1
        return client_id

    def print_match_table(self, filename, matches):
        if not matches:
            print(f"--- No clients found in {filename} ---")
            return

        print(f"REDACTION REPORT: {filename}")
        header = f"| {'Location':<12} | {'Original Name':<25} | {'Redacted As':<15} |"
        sep = f"|{'-'*14}|{'-'*27}|{'-'*17}|"
        print(sep)
        print(header)
        print(sep)
        
        for m in matches:
            orig = (m['orig'][:22] + '..') if len(m['orig']) > 24 else m['orig']
            print(f"| {str(m['loc']):<12} | {orig:<25} | {m['cid']:<15} |")
        print(sep + "\n")

    # --- PPTX LOGIC ---
    def redact_pptx(self, file_path, output_dir):
        prs = Presentation(file_path)
        filename = os.path.basename(file_path)
        modified, file_matches = False, []
        
        for idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not shape.has_text_frame: continue
                for paragraph in shape.text_frame.paragraphs:
                    full_para_text = "".join(run.text for run in paragraph.runs)
                    if not full_para_text.strip(): continue

                    doc = nlp(full_para_text)
                    entities = sorted([e for e in doc.ents if e.label_ == "ORG"], 
                                   key=lambda x: x.start_char, reverse=True)

                    temp_text = full_para_text
                    replaced = False
                    for ent in entities:
                        cid = self.get_client_id(ent.text)
                        if cid:
                            file_matches.append({'loc': f"Slide {idx+1}", 'orig': ent.text, 'cid': cid})
                            temp_text = temp_text[:ent.start_char] + cid + temp_text[ent.end_char:]
                            replaced = modified = True
                    
                    if replaced:
                        for i, run in enumerate(paragraph.runs):
                            run.text = temp_text if i == 0 else ""

        self.print_match_table(filename, file_matches)
        if modified: prs.save(os.path.join(output_dir, f"Redacted_{filename}"))

    # --- EXCEL LOGIC ---
    def redact_xlsx(self, file_path, output_dir):
        wb = openpyxl.load_workbook(file_path)
        filename = os.path.basename(file_path)
        modified, file_matches = False, []

        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str):
                        doc = nlp(cell.value)
                        entities = sorted([e for e in doc.ents if e.label_ == "ORG"], 
                                       key=lambda x: x.start_char, reverse=True)
                        
                        cell_text = cell.value
                        replaced = False
                        for ent in entities:
                            cid = self.get_client_id(ent.text)
                            if cid:
                                file_matches.append({'loc': f"{sheet.title}!{cell.coordinate}", 'orig': ent.text, 'cid': cid})
                                cell_text = cell_text[:ent.start_char] + cid + cell_text[ent.end_char:]
                                replaced = modified = True
                        
                        if replaced:
                            cell.value = cell_text

        self.print_match_table(filename, file_matches)
        if modified: wb.save(os.path.join(output_dir, f"Redacted_{filename}"))

    def save_state(self):
        with open(self.mapping_file, 'w') as f:
            json.dump({'map': self.client_map, 'counter': self.counter}, f, indent=4)

# --- Main Execution ---
input_folder = "./source_data" # Folder containing your PPTX and XLSX
output_folder = "./sanitized_data"
os.makedirs(output_folder, exist_ok=True)
os.makedirs(input_folder, exist_ok=True)

redactor = GlobalRedactor(target_clients=GLOBAL_CLIENT_LIST)

# Get all supported files
files = sorted([f for f in os.listdir(input_folder) if f.endswith(('.pptx', '.xlsx'))])

if not files:
    print(f"No compatible files found in {input_folder}.")
else:
    for file in files:
        full_path = os.path.join(input_folder, file)
        if file.endswith('.pptx'):
            redactor.redact_pptx(full_path, output_folder)
        elif file.endswith('.xlsx'):
            redactor.redact_xlsx(full_path, output_folder)

    redactor.save_state()
    print(f"Success: Processing complete. Mappings updated in client_mapping.json.")