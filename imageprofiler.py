import io
import os
import hashlib
import json
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from PIL import Image
import imagehash

class LogoRedactor:
    def __init__(self, mapping_file="client_mapping.json", clipping_file="logo_clippings.json", clipping_dir="./logo_clippings", threshold=2):
        self.mapping_file = mapping_file
        self.clipping_file = clipping_file
        self.clipping_dir = clipping_dir
        self.threshold = threshold
        self.global_inventory = {}
        
        os.makedirs(self.clipping_dir, exist_ok=True)

        # 1. Load standard client mapping (Text-based)
        if os.path.exists(mapping_file):
            with open(mapping_file, 'r') as f:
                data = json.load(f)
                self.client_map = data.get('map', {})
        else:
            self.client_map = {}

        # 2. Load or initialize logo manifest (Image-based)
        if os.path.exists(clipping_file):
            with open(clipping_file, 'r') as f:
                self.logo_manifest = json.load(f)
                # Determine starting index for new image IDs
                existing_ids = [int(v['id'].split('-')[1]) for v in self.logo_manifest.values() if 'img-' in v.get('id', '')]
                self.img_counter = max(existing_ids) + 1 if existing_ids else 1
        else:
            self.logo_manifest = {}
            self.img_counter = 1

    def get_fingerprints(self, image_blob):
        binary_hash = hashlib.sha256(image_blob).hexdigest()
        img = Image.open(io.BytesIO(image_blob))
        vis_hash = str(imagehash.phash(img))
        return binary_hash, vis_hash, img

    def profile_images(self, file_path):
        """Pass 1: Catalog images and detect existing IDs."""
        prs = Presentation(file_path)
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    b_hash, v_hash, img_obj = self.get_fingerprints(shape.image.blob)
                    
                    if v_hash not in self.global_inventory:
                        clipping_path = os.path.join(self.clipping_dir, f"{v_hash}.png")
                        if not os.path.exists(clipping_path):
                            img_obj.save(clipping_path)

                        self.global_inventory[v_hash] = {
                            "count": 0,
                            "clipping_path": clipping_path,
                            "occurrences": [],
                            "alt_text": getattr(shape, 'alternative_text', None)
                        }
                    
                    self.global_inventory[v_hash]["count"] += 1
                    self.global_inventory[v_hash]["occurrences"].append({
                        "file": file_path,
                        "slide_index": slide_idx,
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height
                    })

    def get_persistent_label(self, v_hash, v_hash_data):
        """Assigns a persistent ID or retrieves an existing one."""
        # A. Check if this hash already has an ID in the saved manifest
        if v_hash in self.logo_manifest:
            return self.logo_manifest[v_hash]["id"]

        # B. Check Alt-Text for a match in client_mapping.json
        alt = (v_hash_data.get("alt_text") or "").lower()
        for client_name, client_id in self.client_map.items():
            if client_name in alt:
                return client_id

        # C. Generate new persistent Image ID
        new_id = f"img-{self.img_counter}"
        self.img_counter += 1
        return new_id

    def process_redaction(self, output_dir):
        """Pass 2: Apply redaction with persistent IDs."""
        logo_hashes = [vh for vh, d in self.global_inventory.items() if d["count"] >= self.threshold]
        
        files_to_process = {}
        for v_hash in logo_hashes:
            data = self.global_inventory[v_hash]
            persistent_id = self.get_persistent_label(v_hash, data)
            
            # Update manifest state
            self.logo_manifest[v_hash] = {
                "id": persistent_id,
                "clipping_path": data["clipping_path"],
                "frequency": data["count"]
            }

            for occ in data["occurrences"]:
                fname = occ["file"]
                if fname not in files_to_process:
                    files_to_process[fname] = []
                occ['display_id'] = persistent_id
                files_to_process[fname].append(occ)

        redacted_total = 0
        for original_path, occurrences in files_to_process.items():
            prs = Presentation(original_path)
            for occ in occurrences:
                slide = prs.slides[occ["slide_index"]]
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        if shape.left == occ["left"] and shape.top == occ["top"]:
                            
                            # Final Label Format: [img-1]-logo
                            display_text = f"[{occ['display_id']}]-logo"
                            
                            rect = slide.shapes.add_shape(1, shape.left, shape.top, shape.width, shape.height)
                            rect.text = display_text
                            
                            # Styling (Centered, Bold, Small)
                            for paragraph in rect.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.size = Pt(8)
                                    run.font.bold = True

                            sp = shape._element
                            sp.getparent().remove(sp)
                            redacted_total += 1
            
            prs.save(os.path.join(output_dir, f"IMG_Sanitized_{os.path.basename(original_path)}"))
        
        # Save manifest for future sessions
        with open(self.clipping_file, 'w') as f:
            json.dump(self.logo_manifest, f, indent=4)
        
        return redacted_total

# --- Execution ---
input_dir = "./source_logo"
output_dir = "./img_sanitized_data"
os.makedirs(output_dir, exist_ok=True)

redactor = LogoRedactor()

files = [f for f in os.listdir(input_dir) if f.endswith(".pptx")]
for f in files:
    redactor.profile_images(os.path.join(input_dir, f))

total = redactor.process_redaction(output_dir)
print(f"Success! {total} logos redacted. Manifest: logo_clippings.json")